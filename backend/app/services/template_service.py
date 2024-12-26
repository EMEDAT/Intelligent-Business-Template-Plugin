import os
import io
import re
import pdfkit
import tempfile
from docx import Document
from pptx import Presentation
from typing import List, Dict

def format_template_content(title: str, key_points: List[str]) -> str:
    """Formats template content with numbering and newlines."""
    if not key_points:
        return f"{title}:\nNo key points provided."

    numbered_points = []
    for i, point in enumerate(key_points):
        numbered_points.append(f"{i + 1}. {point}")

    return f"{title}:\n\n" + "\n".join(numbered_points)

def create_template(template_type: str, key_points: List[str]) -> Dict[str, str]:
    if template_type == "business_plan":
        return generate_business_plan(key_points)
    elif template_type == "pitch_deck":
        return generate_pitch_deck(key_points)
    elif template_type == "marketing_strategy":
        return generate_marketing_strategy(key_points)
    else:
        return {"error": "Unknown template type", "content": ""}

def generate_business_plan(key_points: List[str]) -> Dict[str, str]:
    content = format_template_content("Business Plan", key_points)
    return {"template_type": "business_plan", "content": content}

def generate_pitch_deck(key_points: List[str]) -> Dict[str, str]:
    content = format_template_content("Pitch Deck", key_points)
    return {"template_type": "Pitch Deck", "content": content}

def generate_marketing_strategy(key_points: List[str]) -> Dict[str, str]:
    content = format_template_content("Marketing Strategy", key_points)
    return {"template_type": "Marketing Strategy", "content": content}

def export_as_word(template_content: str) -> bytes:
    if not template_content:
        return b""
    try:
        doc = Document()
        doc.add_paragraph(template_content)
        file_stream = io.BytesIO()
        doc.save(file_stream)
        return file_stream.getvalue()  # Correct: Return the bytes
    except Exception as e:
        print(f"Word Export Error: {e}")
        return b""

def export_as_pdf(template_content: str) -> bytes:
    if not template_content:
        return b""
    try:
        try:
            template_type, rest_of_content = template_content.split(":", 1)
        except ValueError: # Handle cases where there is no colon
            template_type = "Template"
            rest_of_content = template_content
        formatted_content = rest_of_content.replace('\n', '<br>')
        
        html_content = f"<h1><b>{template_type.strip()}</b></h1><p>{formatted_content.strip()}</p>" # bolding the template type
        pdf_bytes = pdfkit.from_string(html_content, False)
        return pdf_bytes
    except Exception as e:
        print(f"PDF Export Error: {e}")
        return b""

def export_as_pptx(template_content: str) -> bytes:
    if not template_content:
        return b""
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tf = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height).text_frame
        tf.text = template_content
        file_stream = io.BytesIO()
        prs.save(file_stream)
        return file_stream.getvalue() # Correct: Return the bytes
    except Exception as e:
        print(f"PowerPoint Export Error: {e}")
        return b""

def handle_temporary_file(template_content: str, file_name: str, format_type: str) -> str:
    with tempfile.TemporaryDirectory() as temp_dir:
        extension = "docx" if format_type == "word" else format_type
        file_path = os.path.join(temp_dir, f"{file_name}.{extension}")
        if format_type == "word":
            export_as_word(template_content, file_path)
        elif format_type == "pdf":
            export_as_pdf(template_content, file_path)
        elif format_type == "pptx":
            export_as_pptx(template_content, file_path)
        else:
            raise ValueError(f"Unsupported format: {format_type}")
        return file_path